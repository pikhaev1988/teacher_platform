#apps/generator/service
import markdown
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Pt
import requests
from django.conf import settings
import io
import uuid
import os
from dotenv import load_dotenv
import uuid
import base64
import urllib3

load_dotenv()
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)
class GPTService:

    AUTH_URL = "https://ngw.devices.sberbank.ru:9443/api/v2/oauth"
    API_URL = "https://gigachat.devices.sberbank.ru/api/v1/chat/completions"

    SYSTEM_PROMPT = """
Ты — профессиональный методист образовательной организации РФ.

Составляй поурочные планы строго в формате ФГОС.
Не добавляй эмодзи.
Не меняй порядок разделов.
Не пропускай разделы.
Пиши официальным методическим языком.

Структура ОБЯЗАТЕЛЬНА:

Ф.И.О. учителя:
Дата:
Класс:
Предмет:
Тема урока:
Тип урока:

Планируемые результаты:

Личностные:
Метапредметные:
Предметные:

Цель урока:

Ключевые слова:

Базовые понятия:

Виды деятельности:

ХОД УРОКА:

ЭТАП 1. МОТИВАЦИОННО-ЦЕЛЕВОЙ
ЭТАП 2. ОПЕРАЦИОННО-ДЕЯТЕЛЬНОСТНЫЙ
ЭТАП 3. РЕФЛЕКСИВНО-ОЦЕНОЧНЫЙ

В конце добавь 2 задания формата ВПР.
"""

    def __init__(self):
        self.client_id = settings.GIGACHAT_CLIENT_ID
        # В .env GIGACHAT_AUTH_KEY уже приходит в том виде, как его даёт Сбер
        # (Authorization key), поэтому его не нужно дополнительно кодировать.
        self.auth_key = settings.GIGACHAT_AUTH_KEY
        self.access_token = self._get_access_token()

    def _get_access_token(self):
        try:
            headers = {
                "Authorization": f"Basic {self.auth_key}",
                "Content-Type": "application/x-www-form-urlencoded",
                "Accept": "application/json",
                "RqUID": str(uuid.uuid4()),
            }

            response = requests.post(
                self.AUTH_URL,
                headers=headers,
                data={"scope": "GIGACHAT_API_PERS"},
                timeout=20,
                verify=False,
            )

            response.raise_for_status()
            data = response.json()

            if "access_token" not in data:
                raise Exception(f"Access token not found: {data}")

            return data["access_token"]

        except requests.Timeout:
            raise Exception("Timeout while getting GigaChat access token")

        except requests.RequestException as e:
            raise Exception(f"GigaChat auth failed: {str(e)}")

        except Exception as e:
            raise Exception(f"Unexpected auth error: {str(e)}")

    def generate_lesson_plan(self, topic, grade, subject):
        try:
            user_prompt = f"""
Сформируй поурочный план.

Класс: {grade}
Предмет: {subject}
Тема урока: {topic}
"""

            headers = {
                "Authorization": f"Bearer {self.access_token}",
                "Content-Type": "application/json"
            }

            payload = {
                "model": "GigaChat",
                "messages": [
                    {"role": "system", "content": self.SYSTEM_PROMPT},
                    {"role": "user", "content": user_prompt}
                ],
                "temperature": 0.3
            }

            response = requests.post(
                self.API_URL,
                headers=headers,
                json=payload,
                timeout=60,
                verify=False,
            )

            response.raise_for_status()
            data = response.json()

            if "choices" not in data:
                return {
                    "success": False,
                    "error": f"Invalid API response: {data}"
                }

            text = data["choices"][0]["message"]["content"]

            return {
                "success": True,
                "text": text,
                "prompt_tokens": data.get("usage", {}).get("prompt_tokens", 0),
                "completion_tokens": data.get("usage", {}).get("completion_tokens", 0),
            }

        except requests.Timeout:
            return {
                "success": False,
                "error": "Timeout while generating lesson plan"
            }

        except requests.RequestException as e:
            return {
                "success": False,
                "error": f"GigaChat request failed: {str(e)}"
            }

        except Exception as e:
            return {
                "success": False,
                "error": f"Unexpected error: {str(e)}"
            }



class DocumentGenerator:
    """Сервис для генерации документов"""

    @staticmethod
    def create_docx_from_markdown(text, topic):
        doc = Document()

        title = doc.add_heading(f'Поурочный план: {topic}', level=1)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        lines = text.split('\n')

        for line in lines:
            if line.startswith('# '):
                doc.add_heading(line[2:], level=1)
            elif line.startswith('## '):
                doc.add_heading(line[3:], level=2)
            elif line.startswith('- ') or line.startswith('* '):
                doc.add_paragraph(line[2:], style='List Bullet')
            elif line.strip():
                p = doc.add_paragraph(line)
                p.paragraph_format.space_after = Pt(6)

        file_stream = io.BytesIO()
        doc.save(file_stream)
        file_stream.seek(0)

        return file_stream

    @staticmethod
    def create_pptx_from_markdown(text, topic):
        prs = Presentation()

        # Титульный слайд
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = topic
        slide.placeholders[1].text = "Поурочный план"

        lines = [l.rstrip() for l in text.split('\n')]
        has_markdown = any(l.startswith('# ') or l.startswith('- ') for l in lines)

        if has_markdown:
            current_title = None
            content = []

            for line in lines:
                if line.startswith('# '):
                    if current_title and content:
                        DocumentGenerator._add_slide(prs, current_title, content)
                    current_title = line[2:]
                    content = []
                elif line.startswith('- '):
                    content.append(line[2:])
                elif line.strip():
                    content.append(line)

            if current_title and content:
                DocumentGenerator._add_slide(prs, current_title, content)
        else:
            # Резервный вариант: обычный текст без markdown.
            # Делим на абзацы по пустым строкам и делаем по слайду на абзац.
            paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]

            for paragraph in paragraphs:
                plines = [l.strip() for l in paragraph.split('\n') if l.strip()]
                if not plines:
                    continue

                first = plines[0]
                # Если это похоже на заголовок этапа – используем как заголовок слайда
                if first.upper().startswith('ЭТАП') or first.endswith(':'):
                    title = first
                    content = plines[1:] or [first]
                else:
                    title = first[:80]
                    content = plines[1:] if len(plines) > 1 else plines

                DocumentGenerator._add_slide(prs, title, content)

        file_stream = io.BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)

        return file_stream

    @staticmethod
    def _add_slide(prs, title, content):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()

        for line in content:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0